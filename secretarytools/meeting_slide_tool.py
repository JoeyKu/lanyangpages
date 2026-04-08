#!/usr/bin/env python3
"""
用法：
    python3 replace.py <file.pptx> --year 2024 --month 7 --speaker 張三 --supervisor 李四
    python3 replace.py <file.pptx> --year 2024 --month 7 --speaker 張三 --supervisor 李四 --report report.pptx

選項：
    --report <report.pptx>  把 report.pptx 第 2 張起全部複製，插入到主檔「工作報告」投影片之後
    --debug                 印出所有含佔位符的 paragraph 原始 XML（診斷用）

佔位符對應：
    [[Title]]     → {year}年{month}月月例會
    [[宣講員]]    → {speaker}
    [[上級指導]]  → {supervisor}

輸出：在原檔名後加 _replaced，例如 template_replaced.pptx
"""

import argparse
import sys
from copy import deepcopy
from pathlib import Path, PurePosixPath

from lxml import etree
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.opc.package import Part, _Relationship, RTM
from pptx.opc.packuri import PackURI
from pptx.parts.image import ImagePart

A_NS    = "http://schemas.openxmlformats.org/drawingml/2006/main"
XML_NS  = "http://www.w3.org/XML/1998/namespace"
T_TAG   = f"{{{A_NS}}}t"
RPR_TAG = f"{{{A_NS}}}rPr"

# ── 文字替換 ──────────────────────────────────────────────────────────────────

def _set_t_text(r_elem, text):
    t = r_elem.find(T_TAG)
    if t is None:
        t = etree.SubElement(r_elem, T_TAG)
    t.text = text
    if text and text != text.strip():
        t.set(f"{{{XML_NS}}}space", "preserve")
    elif f"{{{XML_NS}}}space" in t.attrib:
        del t.attrib[f"{{{XML_NS}}}space"]

def _clone_rPr(src_r):
    rpr = src_r.find(RPR_TAG)
    return deepcopy(rpr) if rpr is not None else None

def _ensure_rPr(r_elem, rpr_clone):
    if rpr_clone is None:
        return
    existing = r_elem.find(RPR_TAG)
    if existing is not None:
        r_elem.remove(existing)
    r_elem.insert(0, rpr_clone)

def replace_in_paragraph(paragraph, mapping, debug=False):
    runs = paragraph.runs
    if not runs:
        return
    full_text = "".join(r.text for r in runs)
    if not any(ph in full_text for ph in mapping):
        return
    if debug:
        print(f"\n[DEBUG] before:\n{etree.tostring(paragraph._p, pretty_print=True).decode()}")

    all_in_single = all(any(ph in r.text for r in runs) for ph in mapping if ph in full_text)
    if all_in_single:
        for run in runs:
            new = run.text
            for ph, val in mapping.items():
                new = new.replace(ph, val)
            if new != run.text:
                _set_t_text(run._r, new)
    else:
        new_text = full_text
        for ph, val in mapping.items():
            new_text = new_text.replace(ph, val)
        def has_sz(run):
            rpr = run._r.find(RPR_TAG)
            return rpr is not None and rpr.get("sz") is not None
        target_run = next((r for r in runs if has_sz(r)), runs[0])
        target_rpr = _clone_rPr(target_run._r)
        for run in runs:
            _set_t_text(run._r, "")
        _set_t_text(runs[0]._r, new_text)
        _ensure_rPr(runs[0]._r, target_rpr)

    if debug:
        print(f"[DEBUG] after:\n{etree.tostring(paragraph._p, pretty_print=True).decode()}")

def process_shape(shape, mapping, debug=False):
    if shape.shape_type == 6:
        for s in shape.shapes:
            process_shape(s, mapping, debug)
        return
    if not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        replace_in_paragraph(para, mapping, debug)


# ── Word 提案解析與插入 ────────────────────────────────────────────────────────

def parse_docx_proposals(doc_path):
    import zipfile
    import xml.etree.ElementTree as ET
    
    items = []
    try:
        with zipfile.ZipFile(str(doc_path), 'r') as z:
            content = z.read('word/document.xml')
            root = ET.fromstring(content)
            namespaces = {'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'}
            
            body = root.find(f".//{{{namespaces['w']}}}body")
            
            paragraphs = []
            if body is not None:
                for child in body:
                    if child.tag == f"{{{namespaces['w']}}}p":
                        text = ''.join(t.text for t in child.findall('.//w:t', namespaces) if t.text)
                        if text.strip():
                            paragraphs.append(text.strip())

            in_section = False
            current_item = None
            
            for text in paragraphs:
                if '宣讀上次決議案執行成效' in text:
                    in_section = True
                    continue
                
                if in_section:
                    if text.startswith('【提案') and '】' in text:
                        if current_item:
                            items.append(current_item)
                        current_item = {'projectNumber': text, 'AA': '', 'BB': ''}
                    elif current_item:
                        stripped_text = text.replace(' ', '').replace('　', '')
                        if stripped_text.startswith('案由：') or stripped_text.startswith('案由:'):
                            sep = '：' if '：' in text else ':'
                            current_item['AA'] = text.split(sep, 1)[1].strip() if sep in text else text
                        elif stripped_text.startswith('執行成效：') or stripped_text.startswith('執行成效:') or stripped_text.startswith('執行辦法：') or stripped_text.startswith('執行辦法:'):
                            sep = '：' if '：' in text else ':'
                            current_item['BB'] = text.split(sep, 1)[1].strip() if sep in text else text
                        elif '工作報告' in text or '提案討論' in text:
                            break
                        elif current_item['BB'] != '':
                            current_item['BB'] += '\n' + text
                        elif current_item['AA'] != '':
                            current_item['AA'] += '\n' + text
            if current_item:
                items.append(current_item)
    except Exception as e:
        print(f"⚠️ 讀取 Word 檔案 {doc_path} 失敗: {e}", file=sys.stderr)
        
    return items

def move_slide_to(prs, slide_to_move, insert_index):
    from pptx.oxml.ns import qn
    prs_elm = prs.part._element
    sldIdLst = prs_elm.find(qn("p:sldIdLst"))
    
    slide_id = slide_to_move.slide_id
    all_sldIds = list(sldIdLst.findall(qn("p:sldId")))
    
    target_el = next((el for el in all_sldIds if int(el.get("id")) == slide_id), None)
    if target_el is not None:
        sldIdLst.remove(target_el)
        all_sldIds_now = list(sldIdLst.findall(qn("p:sldId")))
        if insert_index < len(all_sldIds_now):
            all_sldIds_now[insert_index].addprevious(target_el)
        else:
            sldIdLst.append(target_el)

def insert_proposal_slides(prs, doc_path, debug=False):
    items = parse_docx_proposals(doc_path)
    if not items:
        return
        
    try:
        # 取得母片第 5 張 (index 4)
        layout = prs.slide_layouts[4]
    except IndexError:
        print("⚠️ 找不到母片的第 5 張投影片版面配置，無法插入提案投影片。", file=sys.stderr)
        return
        
    # Find insertion point
    keyword1 = "決議案執行成效"
    keyword2 = "決議執行成效"
    insert_idx = len(prs.slides) - 1
    for i, slide_obj in enumerate(prs.slides):
        text = "\n".join(shape.text for shape in slide_obj.shapes if shape.has_text_frame)
        if keyword1 in text or keyword2 in text:
            insert_idx = i
            break

    print(f"📋 從「{doc_path.name}」建立 {len(items)} 張提案投影片...")
    
    from copy import deepcopy
    template_non_placeholders = []
    mapping_keys = ["{{ProjectNumber}}", "{{project}}", "{{work}}", "{{project_title}}", "{{work_title}}"]
    
    # 預先把非佔位符的模板圖案抽出來，然後唯獨把母片上的變數清除成空字串，以確保「案由」等固定文字會被保留在母片供調整，且不會有疊字問題
    for l_shape in layout.shapes:
        print(l_shape.text)
        if not getattr(l_shape, "is_placeholder", False) and getattr(l_shape, "has_text_frame", False):
           
            if any(k in l_shape.text for k in mapping_keys):
                template_non_placeholders.append(deepcopy(l_shape._element))
                for p in l_shape.text_frame.paragraphs:
                    print("  x" + p.text)
                    full_text = "".join(r.text for r in p.runs)
                    if full_text in mapping_keys:
                        for run in p.runs:
                            print("    y" + run.text)
                            run.text = run.text.replace(run.text, "")
                    #for run in p.runs:
                    #    print("    y" + run.text)
                    #    for k in mapping_keys:
                    #        if k in run.text:
                    #            run.text = run.text.replace(k, "")
                
    for offset, item in enumerate(items):
        slide = prs.slides.add_slide(layout)
        
        mapping = {
            "{{ProjectNumber}}": item['projectNumber'],
            "{{project_title}}": "案        由：",
            "{{project}}": item['AA'],
            "{{work_title}}": "執行成效：",
            "{{work}}": item['BB']
        }

        # 將 layout 中 placeholders 的預設文字完整複製到 slide 上，否則 python-pptx 預設會是空的，導致無法取代
        for l_shape in layout.shapes:
            if getattr(l_shape, "is_placeholder", False) and l_shape.has_text_frame:
                if any(k in l_shape.text for k in mapping_keys):
                    try:
                        s_shape = slide.placeholders[l_shape.placeholder_format.idx]
                        if not s_shape.text.strip():
                            s_shape.text_frame._element.clear_content()
                            for p in l_shape.text_frame._element.p_lst:
                                s_shape.text_frame._element.append(deepcopy(p))
                    except Exception:
                        pass
                        
        # 將預先抽出的非佔位符模板拷貝至此張投影片
        for sp_el in template_non_placeholders:
            slide.shapes._spTree.append(deepcopy(sp_el))

        for shape in slide.shapes:
            process_shape(shape, mapping, debug=debug)
            
        target_idx = insert_idx + 1 + offset
        move_slide_to(prs, slide, target_idx)
        
    print(f"   ✅ 建立提案投影片完成（共 {len(items)} 張）")


# ── 投影片複製插入 ────────────────────────────────────────────────────────────

def _set_rel(part, rId, reltype, target, is_external=False):
    """直接把一個 _Relationship 寫入 part._rels._rels，並指定 rId。"""
    base_uri = part._rels._base_uri
    mode = RTM.EXTERNAL if is_external else RTM.INTERNAL
    part._rels._rels[rId] = _Relationship(base_uri, rId, reltype, mode, target)


def _collect_used_partnames(prs):
    names = set()
    for slide in prs.slides:
        for rel in slide.part.rels.values():
            if not rel.is_external:
                names.add(str(rel.target_part.partname))
    return names


def _unique_partname(existing: set, suffix: str) -> str:
    n = 1
    while True:
        candidate = f"/ppt/media/image{n}{suffix}"
        if candidate not in existing:
            return candidate
        n += 1


def _insert_slide(dest_prs, src_slide, after_index, used_partnames: set):
    """
    複製 src_slide 並插入 dest_prs 的 after_index 之後。

    要點：
    - 用 add_slide 取得合法 sldId；之後換掉 XML 及 rels
    - media part 若 partname 衝突就重新命名（但 rId 保持不變，
      讓 slide XML 裡的 r:id 參照不用修改）
    - 同一個 blob（例如 video 有兩個 rel）只建一個 Part，第二個 rel 重用
    """
    prs_elm  = dest_prs.part._element
    sldIdLst = prs_elm.find(qn("p:sldIdLst"))

    # 1. 選 layout
    src_layout_pn = src_slide.slide_layout.part.partname
    dest_layout = next(
        (l for l in dest_prs.slide_layouts if l.part.partname == src_layout_pn),
        dest_prs.slide_layouts[0]
    )

    # 2. add_slide → 取得有合法 sldId 的新 slide
    new_slide = dest_prs.slides.add_slide(dest_layout)
    new_part  = new_slide.part

    # 3. 換掉 XML（深度複製 src）
    new_part._element = deepcopy(src_slide.part._element)

    # 4. 清掉所有舊 rels，重新建立
    new_part._rels._rels.clear()

    # 記錄「src partname → 已建立的 new Part」，避免同一個 part 建兩次
    built_parts: dict[str, object] = {}

    for rId, rel in src_slide.part.rels.items():
        if rel.is_external:
            _set_rel(new_part, rId, rel.reltype, rel._target, is_external=True)
            continue

        tp = rel.target_part

        # ── slideLayout ──────────────────────────────────────────────────
        if "slideLayout" in rel.reltype:
            _set_rel(new_part, rId, rel.reltype, dest_layout.part)
            continue

        # ── media / image / video ────────────────────────────────────────
        blob = tp._blob
        if not blob:
            continue

        old_pn = str(tp.partname)

        # 同一個 src part 已建過（例如 video 兩個 rel 指向同一個 mp4）
        if old_pn in built_parts:
            _set_rel(new_part, rId, rel.reltype, built_parts[old_pn])
            continue

        suffix = PurePosixPath(old_pn).suffix
        new_pn = old_pn if old_pn not in used_partnames else _unique_partname(used_partnames, suffix)
        used_partnames.add(new_pn)

        if isinstance(tp, ImagePart):
            new_media = ImagePart(PackURI(new_pn), tp.content_type, dest_prs.part.package, blob)
        else:
            new_media = Part(PackURI(new_pn), tp.content_type, dest_prs.part.package, blob)

        built_parts[old_pn] = new_media
        _set_rel(new_part, rId, rel.reltype, new_media)

    # 5. 把 sldId 移到 after_index 之後
    all_sldIds = list(sldIdLst.findall(qn("p:sldId")))
    last_el = all_sldIds[-1]
    sldIdLst.remove(last_el)

    all_sldIds = list(sldIdLst.findall(qn("p:sldId")))
    if after_index + 1 < len(all_sldIds):
        all_sldIds[after_index + 1].addprevious(last_el)
    else:
        sldIdLst.append(last_el)


def _slide_text(slide):
    return "\n".join(
        shape.text_frame.text for shape in slide.shapes if shape.has_text_frame
    )


def insert_report_slides(dest_prs, report_path: Path, keyword="工作報告"):
    report_prs = Presentation(str(report_path))
    total_src  = len(report_prs.slides)
    if total_src < 2:
        print(f"⚠️  {report_path.name} 只有 {total_src} 張，沒有第 2 張可複製。", file=sys.stderr)
        return

    insert_idx = next(
        (i for i, s in enumerate(dest_prs.slides) if keyword in _slide_text(s)), None
    )
    if insert_idx is None:
        print(f"⚠️  找不到含「{keyword}」的投影片，附加到最後。", file=sys.stderr)
        insert_idx = len(dest_prs.slides) - 1

    slides_to_insert = list(report_prs.slides)[1:]
    n = len(slides_to_insert)
    print(f"📋 從「{report_path.name}」複製第 2～{total_src} 張（共 {n} 張），"
          f"插入到第 {insert_idx + 1} 張（「{keyword}」）之後…")

    used_partnames = _collect_used_partnames(dest_prs)

    for offset, src_slide in enumerate(slides_to_insert):
        _insert_slide(dest_prs, src_slide, after_index=insert_idx + offset,
                      used_partnames=used_partnames)

    print(f"   ✅ 插入完成（共 {n} 張）")


def extract_proposal_summary_text(doc_path):
    """
    從 Word 文檔中擷取特定提案文字的總結。
    1. 找到 "宣讀上次決議案執行成效" -> 擷取【提案】、案由、執行成效
    2. 找到 "總會提案討論" 或 "別院提案討論" -> 擷取【提案】、案由，並加上 "執行辦法: "
    3. 遇到 "各類宣導" 停止。
    """
    import zipfile
    import xml.etree.ElementTree as ET
    
    output_lines = []
    try:
        with zipfile.ZipFile(str(doc_path), 'r') as z:
            content = z.read('word/document.xml')
            root = ET.fromstring(content)
            namespaces = {'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'}
            
            body = root.find(f".//{{{namespaces['w']}}}body")
            paragraphs = []
            if body is not None:
                for child in body:
                    if child.tag == f"{{{namespaces['w']}}}p":
                        text = ''.join(t.text for t in child.findall('.//w:t', namespaces) if t.text)
                        if text:
                            paragraphs.append(text.strip())

            # 狀態追蹤
            in_effect_section = False
            # 狀態追蹤
            in_effect_section = False
            in_discussion_section = False
            last_marker = None # 記錄上一個處理的欄位 (aa, bb)
            
            for text in paragraphs:
                # 檢查區域切換
                if '宣讀上次決議案執行成效' in text:
                    in_effect_section = True
                    in_discussion_section = False
                    last_marker = None
                    output_lines.append(f"\n====================宣讀上次決議案執行成效====================")
                    continue
                if '總會提案討論' in text:
                    in_effect_section = False
                    in_discussion_section = True
                    last_marker = None
                    output_lines.append(f"\n====================總會提案討論====================")
                    continue
                if '別院提案討論' in text:
                    in_effect_section = False
                    in_discussion_section = True
                    last_marker = None
                    output_lines.append(f"\n====================別院提案討論====================")
                    continue
                if '各類宣導' in text:
                    in_effect_section = False
                    in_discussion_section = False
                    break
                
                stripped = text.replace(' ', '').replace('　', '')

                # 區域內處理
                if in_effect_section:
                    if text.startswith('【提案') and '】' in text:
                        if last_marker == 'bb':
                            output_lines.append("") # 提案間留空行
                        output_lines.append(text)
                        last_marker = 'pro'
                    elif stripped.startswith('案由：') or stripped.startswith('案由:'):
                        output_lines.append(text)
                        last_marker = 'aa'
                    elif stripped.startswith('執行成效：') or stripped.startswith('執行成效:') or stripped.startswith('執行辦法：') or stripped.startswith('執行辦法:'):
                        output_lines.append(text)
                        last_marker = 'bb'
                    elif last_marker in ['aa', 'bb']:
                        # 延續上一行內容 (多行處理)
                        output_lines.append(text)
                
                elif in_discussion_section:
                    if text.startswith('【提案') and '】' in text:
                        output_lines.append(text)
                        last_marker = 'pro'
                    elif stripped.startswith('案由：') or stripped.startswith('案由:'):
                        output_lines.append(text)
                        output_lines.append("執行辦法: ")
                        output_lines.append("") # 提案間留空行
                        last_marker = 'aa'
                    elif last_marker == 'aa':
                        if any(k in stripped for k in ['說明', '討論', '辦法']):
                            last_marker = 'skip'
                        else:
                            # 案由的多行內容 (在遇到說明、討論、辦法之前)
                            output_lines.insert(-2, text)
                            
    except Exception as e:
        return f"⚠️ 文字擷取失敗: {e}"
        
    return "\n".join(output_lines).strip()


# ── 主流程 ────────────────────────────────────────────────────────────────────

def replace_pptx(input_path, mapping, output_path, report_path=None, doc_path=None, debug=False):
    prs = Presentation(str(input_path))
    for slide in prs.slides:
        for shape in slide.shapes:
            process_shape(shape, mapping, debug)
        if slide.has_notes_slide:
            for shape in slide.notes_slide.shapes:
                process_shape(shape, mapping, debug)
    if report_path is not None:
        insert_report_slides(prs, report_path)
    if doc_path is not None:
        insert_proposal_slides(prs, doc_path, debug)
    prs.save(str(output_path))
    print(f"✅ 已儲存：{output_path}")


# ── CLI ───────────────────────────────────────────────────────────────────────

def run_replace(year, month, speaker, supervisor, has_report, has_doc):
    mapping = {
        "[[Title]]":    f"{year}年{month}月月例會",
        "[[宣講員]]":   speaker,
        "[[上級指導]]": supervisor,
    }
    input_path = Path('/slide-template.pptx')
    output_path = Path('/output.pptx')
    report_path = Path('/report.pptx') if has_report else None
    doc_path = Path('/doc.docx') if has_doc else None
    
    replace_pptx(input_path, mapping, output_path, report_path=report_path, doc_path=doc_path, debug=False)
    
    # 擷取文字回傳給前端
    summary_text = ""
    if has_doc:
        summary_text = extract_proposal_summary_text(doc_path)
        
    return summary_text
